import sys  
import logging  
import io  
import time  
import re  
import json  
from uuid import uuid4  
  
import tiktoken  
import requests  
from flask import Flask, render_template, request, jsonify  
from flask_cors import CORS  
  
from azure.search.documents import SearchClient  
from azure.search.documents.models import QueryType  
from azure.core.credentials import AzureKeyCredential  
from azure.storage.blob import BlobServiceClient  
  
# 如果你還需要文件分片上傳的功能，可保留以下 import  
import fitz                              # PyMuPDF for PDF extraction  
from docx import Document as DocxDocument  # python-docx for .docx/.doc extraction  
from pptx import Presentation as PptxPresentation  # python-pptx for .pptx/.ppt extraction  
  
# -------------------------------------------------------------------  
# 1. 設定 UTF-8 Logger  
logger = logging.getLogger(__name__)  
logger.setLevel(logging.INFO)  
  
utf8_handler = logging.StreamHandler(stream=sys.stdout)  
utf8_handler.setLevel(logging.INFO)  
formatter = logging.Formatter(  
    fmt="%(asctime)s %(levelname)s: %(message)s",  
    datefmt="%Y-%m-%d %H:%M:%S"  
)  
utf8_handler.setFormatter(formatter)  
  
logger.handlers.clear()  
logger.addHandler(utf8_handler)  
  
# -------------------------------------------------------------------  
app = Flask(__name__)  
CORS(app)  
  
# -------------------------------------------------------------------  
# Azure OpenAI 設定  
OAI_ENDPOINT = (  
    "Your_Azure_OpenAi_endpoint"  
)  
OAI_API_KEY = "Your_Azure_OpenAi_API_key" 
OAI_HEADERS = {  
    "Content-Type": "application/json",  
    "api-key": OAI_API_KEY  
}  
  
# -------------------------------------------------------------------  
# Azure Cognitive Search 設定  
SEARCH_ENDPOINT = "Your_Azure_search_service_endpoint"  
SEARCH_API_KEY = "Your_Azure_search_service_API_key"  
INDEX_NAME = "Your_Azure_Indexe_name"  # Replace with your index name 
  
search_client = SearchClient(  
    endpoint=SEARCH_ENDPOINT,  
    index_name=INDEX_NAME,  
    credential=AzureKeyCredential(SEARCH_API_KEY)  
)  
  
# -------------------------------------------------------------------  
# Azure Blob Storage 設定  
STORAGE_CONNECTION_STRING = (  
    "DefaultEndpointsProtocol=https;"  
    "AccountName=onyourdatathiai;"  
    "AccountKey=Y1R3pb5QgNlYHhQ8Og8UteXhOf0xkWRs//QW43yCXjsJOO/"  
    "RfFABPIfEEBm5nhCoo4G+T4l0ws4S+AStVjuBXA==;"  
    "EndpointSuffix=core.windows.net"  
)    
BLOB_CONTAINER_NAME = "Your_Azure_Container_name"               # Replace with your container name  
INDEXER_NAME = "Your_Azure_Indexer_name"          # Replace with your indexer name  
  
blob_service_client = BlobServiceClient.from_connection_string(  
    STORAGE_CONNECTION_STRING  
)  
  
# -------------------------------------------------------------------  
def count_tokens_from_messages(messages, model="gpt-4o-test"):  
    """  
    估算一組 chat messages 的 token 數量（僅用於 log）。  
    """  
    try:  
        encoding = tiktoken.encoding_for_model(model)  
    except KeyError:  
        encoding = tiktoken.get_encoding("cl100k_base")  
    total = 0  
    for msg in messages:  
        total += 4  # <im_start> + role/name  
        for v in msg.values():  
            total += len(encoding.encode(v))  
    total += 2  # <im_end>  
    return total  
  
def chunk_text(text, chunk_size=1000, overlap=200):  
    """  
    將文字按字元數分片（可選擇保留，用於舊版 upload）。  
    """  
    chunks = []  
    start = 0  
    length = len(text)  
    while start < length:  
        end = min(start + chunk_size, length)  
        chunks.append(text[start:end])  
        start += chunk_size - overlap  
    return chunks  
  
def extract_text_from_pdf(stream_bytes):  
    doc = fitz.open(stream=stream_bytes, filetype="pdf")  
    text = ""  
    for page in doc:  
        text += page.get_text()  
    return text  
  
def extract_text_from_docx(stream_bytes):  
    doc = DocxDocument(io.BytesIO(stream_bytes))  
    paragraphs = [p.text for p in doc.paragraphs if p.text]  
    return "\n".join(paragraphs)  
  
def extract_text_from_pptx(stream_bytes):  
    prs = PptxPresentation(io.BytesIO(stream_bytes))  
    text = ""  
    for slide in prs.slides:  
        for shape in slide.shapes:  
            if hasattr(shape, "text") and shape.text:  
                text += shape.text + "\n"  
    return text  
  
# -------------------------------------------------------------------  
@app.route('/')  
def index():  
    return render_template('index.html')  
  
# -------------------------------------------------------------------  
@app.route('/ask_openai', methods=['POST'])  
def ask_openai():  
    try:  
        data = request.json  
        logger.info("[/ask_openai] Received data: %s", data)  
  
        if 'messages' not in data or not isinstance(data['messages'], list):  
            raise ValueError("'messages' field must exist and be a list")  
  
        estimated = count_tokens_from_messages(data['messages'])  
        logger.info("[/ask_openai] Estimated prompt tokens: %d", estimated)  
  
        start = time.time()  
        resp = requests.post(OAI_ENDPOINT, headers=OAI_HEADERS, json=data)  
        elapsed = time.time() - start  
  
        resp.raise_for_status()  
        result = resp.json()  
        logger.info("[/ask_openai] Response: %s", result)  
  
        usage = result.get("usage", {})  
        response_tokens = usage.get("completion_tokens", 0)  
        tps = response_tokens / elapsed if elapsed > 0 else None  
        logger.info("[/ask_openai] completion_tokens=%d, TPS=%.2f", response_tokens, tps)  
  
        result["response_tokens"] = response_tokens  
        result["tps"] = tps  
  
        return jsonify(result)  
    except ValueError as ve:  
        logger.error("[/ask_openai] Invalid data format: %s", ve)  
        return jsonify({"error": str(ve)}), 400  
    except Exception as e:  
        logger.error("[/ask_openai] Error: %s", e)  
        return jsonify({"error": str(e)}), 500  
  
# -------------------------------------------------------------------  
@app.route('/search', methods=['POST'])  
def search():  
    try:  
        data = request.json  
        query = data.get('query', '')  
        results = search_client.search(query, query_type=QueryType.SIMPLE)  
        return jsonify([r for r in results])  
    except Exception as e:  
        logger.error("[/search] Error: %s", e)  
        return jsonify({"error": str(e)}), 500  
  
# -------------------------------------------------------------------  
@app.route('/upload', methods=['POST'])  
def upload():  
    """  
    上傳檔案至 Azure Blob Storage，並觸發 Azure Cognitive Search 的 Indexer。  
    """  
    try:  
        file = request.files['file']  
        overwrite = request.form.get('overwrite') == 'true'  
        filename = file.filename  
  
        blob_client = blob_service_client.get_blob_client(  
            container=BLOB_CONTAINER_NAME,  
            blob=filename  
        )  
  
        logger.info("[/upload] 收到檔案: %s, overwrite=%s", filename, overwrite)  
  
        if not overwrite and blob_client.exists():  
            logger.info("[/upload] 檔案已存在且未選擇覆蓋: %s", filename)  
            return jsonify({"message": "檔案已存在，未覆蓋。"}), 400  
  
        # 上傳檔案到 Blob Storage  
        blob_client.upload_blob(file, overwrite=overwrite)  
        logger.info("[/upload] 成功上傳至 Blob Storage: %s", filename)  
  
        # 觸發 Search Indexer  
        idx_url = f"{SEARCH_ENDPOINT}/indexers/{INDEXER_NAME}/run?api-version=2020-06-30"  
        idx_resp = requests.post(  
            idx_url,  
            headers={  
                "Content-Type": "application/json",  
                "api-key": SEARCH_API_KEY  
            }  
        )  
  
        if idx_resp.status_code == 202:  
            logger.info("[/upload] 成功觸發 Indexer: %s", INDEXER_NAME)  
            return jsonify({"message": "檔案上傳成功，Indexer 已觸發。"}), 200  
        else:  
            logger.error("[/upload] 觸發 Indexer 失敗: %s", idx_resp.text)  
            return jsonify({  
                "message": "檔案上傳成功，但觸發 Indexer 失敗。",  
                "error": idx_resp.text  
            }), 500  
  
    except Exception as e:  
        logger.error("[/upload] 發生錯誤: %s", e)  
        return jsonify({"error": str(e)}), 500  
  
# -------------------------------------------------------------------  
@app.route('/ask', methods=['POST'])  
def ask():  
    try:  
        data = request.json  
        query = data.get('query', '')  
        max_tokens = data.get('max_tokens', 4096)  
        top_p = data.get('top_p', 1.0)  
        temperature = data.get('temperature', 0.7)  
        messages = data.get('messages', [])  
  
        logger.info("[/ask] Received data: %s", data)  
  
        # 1) Retrieval from Cognitive Search  
        search_results = search_client.search(  
            query, top=3, query_type=QueryType.SIMPLE  
        )  
        retrieved = [r.get("content", "") for r in search_results]  
  
        # 2) Prepend retrieved chunks as system messages  
        for chunk in retrieved:  
            messages.append({"role": "system", "content": chunk})  
  
        logger.info("[/ask] final messages payload: %s", json.dumps(messages, ensure_ascii=False, indent=2))  
  
        # 3) Log estimated tokens  
        estimated = count_tokens_from_messages(messages)  
        logger.info("[/ask] Estimated prompt tokens: %d", estimated)  
  
        payload = {  
            "messages": messages,  
            "max_tokens": max_tokens,  
            "top_p": top_p,  
            "temperature": temperature  
        }  
  
        # 4) Call Azure OpenAI  
        start = time.time()  
        resp = requests.post(OAI_ENDPOINT, headers=OAI_HEADERS, json=payload)  
        elapsed = time.time() - start  
        resp.raise_for_status()  
        result = resp.json()  
        logger.info("[/ask] Response: %s", result)  
  
        # 5) 附加 metadata  
        usage = result.get("usage", {})  
        completion_tokens = usage.get("completion_tokens", 0)  
        tps = completion_tokens / elapsed if elapsed > 0 else None  
        result["response_tokens"] = completion_tokens  
        result["tps"] = tps  
  
        return jsonify(result)  
    except Exception as e:  
        logger.error("[/ask] Error: %s", e)  
        return jsonify({"error": str(e)}), 500  
  
# -------------------------------------------------------------------  
if __name__ == '__main__':  
    app.run(debug=True, use_reloader=False)  